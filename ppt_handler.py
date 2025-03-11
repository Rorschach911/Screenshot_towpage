from pptx import Presentation
from pptx.util import Inches, Pt
import os
import win32com.client

def create_ppt_with_screenshots(df, ppt_path, title_text, screenshot_paths):
    """
    创建包含截图和文本的PPT文件，每页显示两组数据
    
    参数:
    df - 包含媒体信息的DataFrame
    ppt_path - PPT文件保存路径
    title_text - 页面标题
    screenshot_paths - 截图文件路径列表
    """
    prs = Presentation()
    
    # 保存PPT以便后续使用COM接口设置默认字体
    prs.save(ppt_path)
    
    # 使用COM接口设置默认字体大小为16
    set_default_font_size(ppt_path)
    
    # 重新打开保存的PPT
    prs = Presentation(ppt_path)
    
    # 每页显示两组数据，计算需要的页数
    total_items = len(df)
    for i in range(0, total_items, 2):
        # 创建新的PPT幻灯片
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
        
        # 左侧数据 - 第一组
        if i < total_items:
            row1 = df.iloc[i]
            # 左侧文本框位置参数
            left1 = Inches(0.5)
            width1 = Inches(4)
            height = Inches(0.5)
            
            # 添加媒体名称
            txBox = slide.shapes.add_textbox(left1, Inches(0.5), width1, height)
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.font.size = Pt(16)
            p.text = f"媒体：{row1['媒体名称']}"
            
            # 添加发布时间
            txBox = slide.shapes.add_textbox(left1, Inches(1), width1, height)
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.font.size = Pt(16)
            p.text = f"发布时间：{row1['发布时间']}"
            
            # 添加页面标题
            txBox = slide.shapes.add_textbox(left1, Inches(1.5), width1, height)
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.font.size = Pt(16)
            p.text = f"标题：{title_text}"
            
            # 添加链接
            txBox = slide.shapes.add_textbox(left1, Inches(2), width1, height)
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.font.size = Pt(16)
            p.text = f"链接：{row1['链接']}"
            
            # 添加截图（缩小到四分之一大小）
            if i < len(screenshot_paths) and os.path.exists(screenshot_paths[i]):
                img_width = Inches(4.5)  # 原始宽度的四分之一
                img_height = Inches(3.5)  # 保持宽高比
                slide.shapes.add_picture(screenshot_paths[i], left1, Inches(2.5), 
                                      width=img_width, height=img_height)
        
        # 右侧数据 - 第二组
        if i + 1 < total_items:
            row2 = df.iloc[i + 1]
            # 右侧文本框位置参数
            left2 = Inches(5.5)
            width2 = Inches(4)
            
            # 添加媒体名称
            txBox = slide.shapes.add_textbox(left2, Inches(0.5), width2, height)
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.font.size = Pt(16)
            p.text = f"媒体：{row2['媒体名称']}"
            
            # 添加发布时间
            txBox = slide.shapes.add_textbox(left2, Inches(1), width2, height)
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.font.size = Pt(16)
            p.text = f"发布时间：{row2['发布时间']}"
            
            # 添加页面标题
            txBox = slide.shapes.add_textbox(left2, Inches(1.5), width2, height)
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.font.size = Pt(16)
            p.text = f"标题：{title_text}"
            
            # 添加链接
            txBox = slide.shapes.add_textbox(left2, Inches(2), width2, height)
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.font.size = Pt(16)
            p.text = f"链接：{row2['链接']}"
            
            # 添加截图（缩小到四分之一大小）
            if i + 1 < len(screenshot_paths) and os.path.exists(screenshot_paths[i + 1]):
                img_width = Inches(4.5)  # 原始宽度的四分之一
                img_height = Inches(3.5)  # 保持宽高比
                slide.shapes.add_picture(screenshot_paths[i + 1], left2, Inches(2.5), 
                                      width=img_width, height=img_height)
    
    # 保存PPT
    prs.save(ppt_path)
    return True

def add_slide_with_screenshot(prs, screenshot_path, media_name, publish_time, title_text, link):
    """
    向现有PPT添加包含截图和文本的幻灯片
    
    参数:
    prs - 现有的Presentation对象
    screenshot_path - 截图文件路径
    media_name - 媒体名称
    publish_time - 发布时间
    title_text - 页面标题
    link - 链接
    
    返回:
    添加的幻灯片对象
    """
    # 创建新的PPT幻灯片
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
    
    # 设置文本框和图片的位置参数
    left = Inches(1)
    width = Inches(8)
    height = Inches(0.5)
    
    # 添加媒体名称
    txBox = slide.shapes.add_textbox(left, Inches(0.5), width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"媒体：{media_name}"
    p.font.size = Pt(16)  # 设置字号为16
    
    # 添加发布时间
    txBox = slide.shapes.add_textbox(left, Inches(1), width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"发布时间：{publish_time}"
    p.font.size = Pt(16)  # 设置字号为16
    
    # 添加页面标题
    txBox = slide.shapes.add_textbox(left, Inches(1.5), width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"标题：{title_text}"
    p.font.size = Pt(16)  # 设置字号为16
    
    # 添加链接
    txBox = slide.shapes.add_textbox(left, Inches(2), width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"链接：{link}"
    p.font.size = Pt(16)  # 设置字号为16
    
    # 添加截图（缩小到二分之一大小）
    if os.path.exists(screenshot_path):
        img_width = Inches(10/2)  # 原始宽度的二分之一
        img_height = Inches(7.5/2)  # 保持宽高比
        slide.shapes.add_picture(screenshot_path, left, Inches(2.5), 
                              width=img_width, height=img_height)
    
    return slide

def create_new_ppt():
    """
    创建一个新的PPT文件，并设置默认字体大小为16
    
    返回:
    新创建的Presentation对象
    """
    prs = Presentation()
    return prs

def save_ppt(prs, ppt_path):
    """
    保存PPT文件
    
    参数:
    prs - Presentation对象
    ppt_path - 保存路径
    
    返回:
    是否保存成功
    """
    try:
        prs.save(ppt_path)
        return True
    except Exception as e:
        print(f"保存PPT时出错: {str(e)}")
        return False
def add_slide_and_content(media_name, publish_time, title_text, link):
    """
    在当前打开的PPT中添加新幻灯片并插入内容
    """
    try:
        # 获取当前打开的PPT应用
        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        presentation = powerpoint.ActivePresentation
        
        # 添加新幻灯片
        slide = presentation.Slides.Add(presentation.Slides.Count + 1, 12)  # 12是空白布局
        
        # 设置内容
        left = 30  # 0.5英寸
        top = 85    # 0.5英寸
        width = 450  # 4英寸
        height = 40  # 0.5英寸
        
        # 创建文本框并设置内容
        textbox = slide.Shapes.AddTextbox(1, left, top, width, height * 4)
        text_frame = textbox.TextFrame
        text_frame.TextRange.Text = (
            f'媒体：{media_name}\n'
            f'发布时间：{publish_time}\n'
            f'标题：{title_text}\n'
            f'链接：{link}'
        )
        text_frame.TextRange.Font.Size = 16
        text_frame.WordWrap = True
        
        return slide
        
    except Exception as e:
        print(f"创建新幻灯片时出错: {str(e)}")
        return None
def add_slide_and_content_pair(media_name1, publish_time1, title_text1, link1, 
                              media_name2=None, publish_time2=None, title_text2=None, link2=None):
    """
    在当前打开的PPT中添加新幻灯片并插入两组内容
    """
    try:
        # 获取当前打开的PPT应用
        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        presentation = powerpoint.ActivePresentation
        
        # 添加新幻灯片
        slide = presentation.Slides.Add(presentation.Slides.Count + 1, 12)  # 12是空白布局
        
        # 左侧内容 - 第一组
        left1 = 30  # 0.5英寸
        top = 85    # 0.5英寸
        width1 = 450  # 4英寸
        height = 40  # 0.5英寸
        
        # 创建左侧文本框并设置内容
        textbox1 = slide.Shapes.AddTextbox(1, left1, top, width1, height * 4)
        text_frame1 = textbox1.TextFrame
        text_frame1.TextRange.Text = (
            f'媒体：{media_name1}\n'
            f'发布时间：{publish_time1}\n'
            f'标题：{title_text1}\n'
            f'链接：{link1}'
        )
        text_frame1.TextRange.Font.Size = 16
        text_frame1.WordWrap = True
        
        # 如果有第二组数据，添加右侧内容
        if media_name2:
            left2 = 490  # 5.5英寸
            width2 = 450  # 4英寸
            
            # 创建右侧文本框并设置内容
            textbox2 = slide.Shapes.AddTextbox(1, left2, top, width2, height * 4)
            text_frame2 = textbox2.TextFrame
            text_frame2.TextRange.Text = (
                f'媒体：{media_name2}\n'
                f'发布时间：{publish_time2}\n'
                f'标题：{title_text2}\n'
                f'链接：{link2}'
            )
            text_frame2.TextRange.Font.Size = 16
            text_frame2.WordWrap = True
        
        return slide
        
    except Exception as e:
        print(f"创建新幻灯片时出错: {str(e)}")
        return None

def set_default_font_size(ppt_path):
    """
    使用COM接口设置PPT的默认字体大小为16
    
    参数:
    ppt_path - PPT文件路径
    """
    try:
        # 确保文件路径是绝对路径
        ppt_path = os.path.abspath(ppt_path)
        
        # 检查文件是否存在
        if not os.path.exists(ppt_path):
            print(f"错误: PPT文件不存在: {ppt_path}")
            return False
            
        # 创建PowerPoint应用程序实例
        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        powerpoint.Visible = True
        
        # 打开PPT文件
        presentation = powerpoint.Presentations.Open(ppt_path)
        
        # 设置默认字体大小
        for design in presentation.Designs:
            for master in design.SlideMaster.CustomLayouts:
                for shape in master.Shapes:
                    if shape.HasTextFrame:
                        shape.TextFrame.TextRange.Font.Size = 16
        
        return presentation
        
    except Exception as e:
        print(f"设置默认字体大小时出错: {str(e)}")
        return None
def open_and_set_font_size(ppt_path):
    """
    打开现有PPT并设置默认字体大小为16
    
    参数:
    ppt_path - PPT文件路径
    """
    return set_default_font_size(ppt_path)